import sys
import os
import json
import spacy
import fitz  # PyMuPDF for PDF
import pytesseract
from PIL import Image
from pptx import Presentation
import pandas as pd
import sqlite3
from transformers import pipeline  
from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer 
import docx
import warnings
warnings.filterwarnings("ignore", category=FutureWarning, module='huggingface_hub')

nlp = spacy.load("en_core_web_sm")


summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
sentiment_analyzer = pipeline("sentiment-analysis", model="distilbert-base-uncased-finetuned-sst-2-english")

vader_analyzer = SentimentIntensityAnalyzer()
inappropriate_words = [  # General Negative Terms
    "damn", "awful", "terrible", "horrible", "dreadful", "poor", "miserable", "unpleasant",
    "disappointing", "sad", "depressing", "distressing", "negative", "unfortunate", "unhappy",
    "dismal", "painful", "distasteful", "regretful", "frustrating", "angry", "upset",
    "annoying", "irritating", "disturbing", "revolting", "repulsive", "sickening", "disgusting",
    "shocking", "appalling", "grotesque", "repugnant", "vile", "loathsome", "wretched",
    "abominable", "detestable", "repellent", "unbearable", "intolerable", "unacceptable",
    "unseemly", "disgraceful", "dismal", "bleak", "regrettable", "heartbreaking",
    "disheartening", "deplorable", "shameful", "atrocious", "heinous", "abominable",

    # Racial and Ethnic Slurs
    "nigger", "chink", "spic", "gook", "raghead", "wetback", "sandnigger", "coon",
    "dago", "guinea", "haji", "honky", "jap", "kraut", "mick", "paki", "pollock",
    "slope", "wop", "yid", "zipperhead", "boong", "gook", "hymie", "kike", "nazi",
    "redskin", "skid", "spade", "towelhead", "tribal", "wetback", "brownie", "spade",

    # Gender and Sexuality Slurs
    "tranny", "trannie", "sexchange", "faggotry", "lesbo", "dyke", "queer", "homosexual",
    "transvestite", "transsexual", "twink", "shemale", "butch", "pansy", "chick", "bender",
    "straight", "gay", "bi", "homo", "sexist", "misogynistic", "misandrist", "chauvinist",
    "rape", "molestation", "abuse", "pervert", "deviant", "slut", "whore", "hooker",
    "escort", "callgirl", "prostitute", "stripper", "stipper", "naked", "bareback",
    "peep", "streaking", "swinger", "incest", "fistfuck", "sadism", "masochism",
    "anal", "pussy", "cum", "jizz", "bukkake", "orgasm", "porn", "sodomy",

    # Discriminatory and Insulting Terms
    "racism", "bigot", "xenophobic", "hate", "bigotry", "discrimination", "abusive",
    "harassment", "bullying", "slur", "discriminatory", "offensive", "insulting",
    "derogatory", "unethical", "profane", "vulgar", "obscene", "indecent", "disgraceful",
    "unseemly", "unfit", "distasteful", "unacceptable", "intolerable", "dismal",
    "grating", "annoying", "irritating", "disturbing", "off-putting", "revolting",
    "repulsive", "despicable", "heinous", "atrocious", "abominable", "wretched",
    "repugnant", "loathsome", "deplorable", "sickening", "shameful", "disgusting",
    "unpleasant", "heartbreaking", "shocking", "appalling", "grotesque", "repellent",
    "revulsive", "reprehensible", "degrading", "demeaning", "oppressive", "threatening",
    "menacing", "harsh", "cruel", "inhuman",

        # General Offensive Terms
    "abuse", "asshole", "bastard", "bitch", "cunt", "dick", "fag", "faggot",
    "fuck", "motherfucker", "piss", "shit", "slut", "twat", "whore", "wanker",
    "douchebag", "cocksucker", "prick", "cocks", "vagina", "penis", "fist", "rape",
    "molest", "pedophile", "retard", "idiot", "moron", "simpleton", "dumbass",
    "dumbfuck", "shithead", "fuckface", "fucktard", "prickface", "shitfuck",
    "slutbag", "whorebag", "cumslut", "pussylips", "ballbag", "cockbag", "cuntface",
    "dickhead", "dickwad", "cocksucking", "asslicking", "pussy", "butt", "asshole",
    "cunt", "cocksucker", "motherfucking", "bastard", "douche", "blowjob", "handjob",
    "cumshot", "deepthroat", "gangbang", "fisting", "sodomize", "sex", "porn", "fetish",
    "anal", "vulgar", "obscene", "indecent", "blasphemy", "sacrilegious", "inappropriate",

    # Racial and Ethnic Slurs
    "nigger", "chink", "spic", "gook", "raghead", "wetback", "cameltoe",
    "sandnigger", "coon", "dago", "guinea", "haji", "honky", "jap", "kraut",
    "mick", "paki", "pollock", "slope", "wop", "yid", "zipperhead", "boong",
    "gook", "hymie", "kike", "nazi", "redskin", "skid", "spade", "towelhead",
    "tribal", "wetback", "brownie", "spade", "racial", "ethnic", "savage",


    "tranny", "trannie", "sexchange", "faggotry", "lesbo", "dyke", "queer", "homosexual",
    "transvestite", "transsexual", "twink", "shemale", "butch", "pansy", "chick",
    "bender", "straight", "gay", "bi", "homo", "sexist", "misogynistic", "misandrist",
    "chauvinist", "rape", "molestation", "abuse", "pervert", "deviant", "slut", "whore",
    "hooker", "escort", "callgirl", "prostitute", "stripper", "stipper", "naked",
    "bareback", "peep", "streaking", "swinger", "incest", "fistfuck", "sadism",
    "masochism", "anal", "pussy", "cum", "jizz", "bukkake", "orgasm", "porn", "sodomy",


    "racism", "bigot", "xenophobic", "hate", "bigotry", "discrimination", "abusive",
    "harassment", "bullying", "slur", "discriminatory", "offensive", "insulting",
    "derogatory", "unethical", "profane", "vulgar", "obscene", "indecent", "disgraceful",
    "unseemly", "unfit", "distasteful", "unacceptable", "intolerable", "dismal", "grating",
    "annoying", "irritating", "disturbing", "off-putting", "revolting", "repulsive",
    "despicable", "heinous", "atrocious", "abominable", "wretched", "repugnant",
    "loathsome", "deplorable", "shameful", "sickening", "disgusting", "unpleasant",
    "heartbreaking", "shocking", "appalling", "grotesque", "repellent", "repugnant",
    "vile", "loathsome", "wretched", "abominable", "distasteful", "insufferable",
    "intolerable", "unbearable", "detestable", "revulsive", "reprehensible",
    "degrading", "degrading", "disrespectful", "insulting", "dehumanizing",
    "demeaning", "oppressive", "threatening", "menacing", "harsh", "cruel", "inhuman"
  
]

def analyze_document(content, chunk_size=1024):
    sentiments = []
    for i in range(0, len(content), chunk_size):
        chunk = content[i:i+chunk_size]
        vader_result = vader_analyzer.polarity_scores(chunk)
        if vader_result['compound'] > 0.05:
            sentiments.append("positive")
        elif vader_result['compound'] < -0.05:
            sentiments.append("negative")
        else:
            sentiments.append("neutral")

    if "inappropriate" in sentiments:
        return "inappropriate"
    elif "negative" in sentiments:
        return "negative"
    elif "positive" in sentiments:
        return "positive"
    else:
        return "neutral"

def generate_summary(content, chunk_size=1024):
    summaries = []
    for i in range(0, len(content), chunk_size):
        chunk = content[i:i+chunk_size]
        summary = summarizer(chunk, max_length=150, min_length=50, do_sample=False)
        summaries.append(summary[0]['summary_text'])
    return " ".join(summaries)

def generate_title(content):
    summary = summarizer(content[:512], max_length=50, min_length=10, do_sample=False)
    title = summary[0]['summary_text'] if summary else "Title not available"
    if title == "Title not available":
        doc = nlp(content)
        first_sentence = list(doc.sents)[0].text if len(list(doc.sents)) > 0 else "Title not available"
        return first_sentence[:50]  
    return title

def extract_text_from_pdf(file_path):
    text = ""
    with fitz.open(file_path) as pdf:
        for page in pdf:
            text += page.get_text()
            if not text.strip():
             
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text += pytesseract.image_to_string(img)
    return text

def extract_text_from_pptx(file_path):
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
            elif shape.shape_type == 13:  
                img = shape.image  
                img_bytes = img.blob  
                with open("temp_image.png", "wb") as f:
                    f.write(img_bytes)  
                img_for_ocr = Image.open("temp_image.png")  
                text += pytesseract.image_to_string(img_for_ocr)  
                os.remove("temp_image.png")  
    return text


def extract_text_from_excel(file_path):
    text = ""
    df = pd.read_excel(file_path, sheet_name=None)
    for sheet_name, sheet_df in df.items():
        text += sheet_df.to_string()
    return text

def extract_text_from_csv(file_path):
    text = ""
    df = pd.read_csv(file_path)
    text += df.to_string()
    return text

def extract_text_from_sql(file_path):
    text = ""
    conn = sqlite3.connect(file_path)
    cursor = conn.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
    tables = cursor.fetchall()
    for table_name in tables:
        cursor.execute(f"SELECT * FROM {table_name[0]}")
        rows = cursor.fetchall()
        for row in rows:
            text += " ".join(map(str, row)) + "\n"
    conn.close()
    return text


def extract_text_from_docx(file_path):
    text = ""
    doc = docx.Document(file_path)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def process_document(file_path):
    base_dir = "priv/static/assets/documents/"
    full_path = os.path.join(base_dir, file_path)
    
    try:
        if not os.path.exists(full_path):
            raise FileNotFoundError(f"File {file_path} does not exist.")

        ext = os.path.splitext(file_path)[1].lower()
        if ext == ".pdf":
            content = extract_text_from_pdf(full_path)
        elif ext == ".pptx" or ext == ".ppt":
            content = extract_text_from_pptx(full_path)
        elif ext == ".xls" or ext == ".xlsx":
            content = extract_text_from_excel(full_path)
        elif ext == ".csv":
            content = extract_text_from_csv(full_path)
        elif ext == ".db": 
            content = extract_text_from_sql(full_path)
        elif ext == ".docx":
            content = extract_text_from_docx(full_path)
        else:
            raise ValueError("Unsupported file type")

        sentiment = analyze_document(content)

        if sentiment == "inappropriate":
            title = "Document Not Processed"
            summary = "This document contains inappropriate content and cannot be processed."
        else:
            summary = generate_summary(content)
            title = generate_title(content)
        
        result = {
            "title": title,
            "summary": summary,
            "sentiment": sentiment
        }
        
        print(json.dumps(result))

    except Exception as e:
        print(json.dumps({"error": str(e)}))

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(json.dumps({"error": "Please provide the file name as an argument."}))
    else:
        file_name = sys.argv[1]
        process_document(file_name)
