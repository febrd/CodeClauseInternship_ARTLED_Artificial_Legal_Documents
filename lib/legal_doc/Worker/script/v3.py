import sys
import os
import json
import spacy
import fitz  
import pytesseract
from PIL import Image
from pptx import Presentation
import pandas as pd
import sqlite3
from transformers import pipeline
from vaderSentiment.vaderSentiment import SentimentIntensityAnalyzer
import docx
import warnings
from tqdm import tqdm  
import magic 
import math  

warnings.filterwarnings("ignore", category=FutureWarning, module='huggingface_hub')

nlp = spacy.load("en_core_web_sm")
summarizer = pipeline("summarization", model="facebook/bart-large-cnn")
sentiment_analyzer = pipeline("sentiment-analysis", model="distilbert-base-uncased-finetuned-sst-2-english")
vader_analyzer = SentimentIntensityAnalyzer()

inappropriate_words = ["badword1", "inappropriate2", "offensive3"]

def analyze_document(content, chunk_size=1024):
    """Analyze sentiment of document content in chunks for large files."""
    sentiments = []
    for i in range(0, len(content), chunk_size):
        chunk = content[i:i + chunk_size]
        vader_result = vader_analyzer.polarity_scores(chunk)
        if vader_result['compound'] > 0.05:
            sentiments.append("positive")
        elif vader_result['compound'] < -0.05:
            sentiments.append("negative")
        else:
            sentiments.append("neutral")

    if any(word in content for word in inappropriate_words):
        return "inappropriate"
    elif "negative" in sentiments:
        return "negative"
    elif "positive" in sentiments:
        return "positive"
    else:
        return "neutral"

def generate_summary2(content, chunk_size=1024, max_length=150, min_length=50):
    """Generate a concise summary dynamically based on content chunks for large files."""
    summaries = []
    num_chunks = math.ceil(len(content) / chunk_size)
    combined_text = ""
    for i in tqdm(range(0, len(content), chunk_size), desc="Processing"):
        chunk = content[i:i + chunk_size]
        combined_text += chunk + " "
    try:
        summary = summarizer(combined_text, max_length=max_length, min_length=min_length, do_sample=False)
        return summary[0]['summary_text']
    except Exception as e:
        print(f"Error during summarization: {e}")
        return combined_text[:150]

def generate_summary(content, chunk_size=1024):
    """Generate summary dynamically based on content chunks for large files."""
    summaries = []
    num_chunks = math.ceil(len(content) / chunk_size)
    for i in tqdm(range(0, len(content), chunk_size), desc="Summarizing"):
        chunk = content[i:i + chunk_size]
        try:
            summary = summarizer(chunk, max_length=150, min_length=50, do_sample=False)
            summaries.append(summary[0]['summary_text'])
        except Exception as e:
            print(f"Error during summarization: {e}")
            summaries.append(chunk[:150])  # Fallback to truncating the chunk if summarization fails
    return " ".join(summaries)

def generate_title(content):
    """Generate a title for the document based on the content."""
    summary = summarizer(content[:512], max_length=50, min_length=10, do_sample=False)
    title = summary[0]['summary_text'] if summary else "Title not available"
    if title == "Title not available":
        doc = nlp(content)
        first_sentence = list(doc.sents)[0].text if len(list(doc.sents)) > 0 else "Title not available"
        return first_sentence[:50]
    return title

def extract_text_from_pdf(file_path):
    """Extract text from PDF files, including OCR for images."""
    text = ""
    with fitz.open(file_path) as pdf:
        for page in pdf:
            text += page.get_text()
            if not text.strip(): 
                pix = page.get_pixmap()
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                text += pytesseract.image_to_string(img)
    return text

def extract_text_from_pptx(file_path):
    """Extract text from PPTX slides, including OCR for embedded images."""
    text = ""
    presentation = Presentation(file_path)
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text
            elif shape.shape_type == 13:  
                img = shape.image
                img_bytes = img.blob
                with open("temp_image.png", "wb") as f:
                    f.write(img_bytes)
                img_for_ocr = Image.open("temp_image.png")
                text += pytesseract.image_to_string(img_for_ocr)
                os.remove("temp_image.png")
    return text

def extract_text_from_excel(file_path):
    """Extract text from Excel files, reading all sheets."""
    text = ""
    df = pd.read_excel(file_path, sheet_name=None)
    for sheet_name, sheet_df in df.items():
        text += sheet_df.to_string()
    return text

def extract_text_from_csv(file_path):
    """Extract text from CSV files."""
    text = ""
    df = pd.read_csv(file_path)
    text += df.to_string()
    return text

def extract_text_from_sql(file_path):
    """Extract text from SQLite databases."""
    text = ""
    conn = sqlite3.connect(file_path)
    cursor = conn.cursor()
    cursor.execute("SELECT name FROM sqlite_master WHERE type='table'")
    tables = cursor.fetchall()
    for table_name in tables:
        cursor.execute(f"SELECT * FROM {table_name[0]}")
        rows = cursor.fetchall()
        for row in rows:
            text += " ".join(map(str, row)) + "\n"
    conn.close()
    return text

def extract_text_from_docx(file_path):
    """Extract text from Word (DOCX) files."""
    text = ""
    doc = docx.Document(file_path)
    for paragraph in doc.paragraphs:
        text += paragraph.text + "\n"
    return text

def detect_file_type(file_path):
    """Dynamically detect file type using libmagic."""
    mime = magic.Magic(mime=True)
    file_type = mime.from_file(file_path)
    return file_type

def process_document(file_path):
    """Process document dynamically based on file type."""
    base_dir = "priv/static/assets/documents/"
    full_path = os.path.join(base_dir, file_path)

    try:
        if not os.path.exists(full_path):
            raise FileNotFoundError(f"File {file_path} does not exist.")

        file_type = detect_file_type(full_path)
        print(f"Detected file type: {file_type}")

        if "pdf" in file_type:
            content = extract_text_from_pdf(full_path)
        elif "vnd.openxmlformats-officedocument.presentationml.presentation" in file_type:
            content = extract_text_from_pptx(full_path)
        elif "vnd.openxmlformats-officedocument.spreadsheetml.sheet" in file_type:
            content = extract_text_from_excel(full_path)
        elif "csv" in file_type:
            content = extract_text_from_csv(full_path)
        elif "sqlite" in file_type:
            content = extract_text_from_sql(full_path)
        elif "msword" in file_type or "officedocument.wordprocessingml" in file_type:
            content = extract_text_from_docx(full_path)
        else:
            raise ValueError("Unsupported file type")

        sentiment = analyze_document(content)

        if sentiment == "inappropriate":
            title = "Document Not Processed"
            summary = "This document contains inappropriate content and cannot be processed."
        else:
            summary = generate_summary(content)
            title = generate_title(content)

        result = {
            "title": title,
            "summary": summary,
            "sentiment": sentiment
        }

        print(json.dumps(result))

    except Exception as e:
        print(json.dumps({"error": str(e)}))

if __name__ == "__main__":
    if len(sys.argv) != 2:
        print(json.dumps({"error": "Please provide the file name as an argument."}))
    else:
        file_name = sys.argv[1]
        process_document(file_name)
